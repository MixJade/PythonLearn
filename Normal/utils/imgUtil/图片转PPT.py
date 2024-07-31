# coding=utf-8
# @Time    : 2024/7/31 21:54
# @Software: PyCharm
import os

from pptx import Presentation
from pptx.util import Cm

# 创建一个PPT对象
prs = Presentation()
prs.slide_width = Cm(19.05)
prs.slide_height = Cm(25.4)

# 获取图片文件夹下所有图片
image_path = "C:\\MyHide\\picacg\\哔咔\\美咲\\"
images = os.listdir(image_path)

# 遍历所有图片文件，然后单独建立一页PPT然后添加图片
for image in images:
    # 添加一页ppt
    slide_layout = prs.slide_layouts[6]  # 使用空白页布局
    slide = prs.slides.add_slide(slide_layout)

    # 设置背景
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.rgb = (255, 255, 255)  # 指定背景为白色

    # 一次性添加一张图片
    slide.shapes.add_picture(image_path + image, 0, 0,
                             width=prs.slide_width,
                             height=prs.slide_height)

# 保存ppt
prs.save('test.pptx')
print('生成成功')
