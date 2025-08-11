# coding=utf-8
# @Time    : 2024/7/31 21:54
# @Software: PyCharm
import os
import re

from PIL import Image
from pptx import Presentation
from pptx.util import Cm, Length


def try_int(s: str) -> int | str:
    try:
        return int(s)
    except ValueError:
        return s.lower()


def alphanum_key(s: str) -> list[int]:
    return [try_int(c) for c in re.split('([0-9]+)', s)]


def get_image_size(directory1: str, image_name1: str) -> tuple[Length, Length, Length, Length]:
    """重设图片的大小，与PPT的空隙
    """
    # 读取文件宽高
    img = Image.open(os.path.join(directory1, image_name1))
    width, height = img.size
    # 查看比例是否正常
    ratio = width / height
    # 按照比例裁剪图片
    if ratio < ppt_ratio:
        # print(f"{image_name1}比例{width, height}小于PPT:高度过大")
        new_ppt_width = ratio * ppt_height
        ppt_left = (ppt_width - new_ppt_width) / 2
        return Cm(new_ppt_width), prs.slide_height, Cm(ppt_left), Cm(0)
    elif ratio > ppt_ratio:
        # print(f"{image_name1}比例{width, height}大于PPT:宽度过大")
        new_ppt_height = ppt_width / ratio
        ppt_top = (ppt_height - new_ppt_height) / 2
        return prs.slide_width, Cm(new_ppt_height), Cm(0), Cm(ppt_top)
    else:
        return prs.slide_width, prs.slide_height, Cm(0), Cm(0)


if __name__ == '__main__':
    # ppt的宽高(A4尺寸)
    ppt_width, ppt_height = 19.05, 25.4
    print(f"图片转PPT,尺寸为{ppt_width, ppt_height}")
    # 获取图片文件夹下所有图片
    file_path = input("输入目标文件夹下某个图片路径：").strip('"')
    # 创建一个PPT对象
    prs = Presentation()
    # ppt的比例
    ppt_ratio = ppt_width / ppt_height
    # ppt页面的大小
    prs.slide_width = Cm(ppt_width)
    prs.slide_height = Cm(ppt_height)
    # 获取文件路径
    directory = os.path.dirname(file_path)
    # 获取文件夹中所有的文件（只读取文件，不读取文件夹）
    image_names = [f for f in os.listdir(directory) if os.path.isfile(os.path.join(directory, f))]
    # 解析字符串中的数字、不区分字母大小写排序
    image_names.sort(key=alphanum_key)
    # 遍历所有图片文件，然后单独建立一页PPT然后添加图片
    for image_name in image_names:
        img_path = os.path.join(directory, image_name)
        # 设置图片宽高
        img_weight, img_height, left_pd, top_pd = get_image_size(directory, image_name)
        # 添加一页ppt
        slide_layout = prs.slide_layouts[6]  # 使用空白页布局
        slide = prs.slides.add_slide(slide_layout)

        # 设置背景
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.rgb = (255, 255, 255)  # 指定背景为白色

        # 一次性添加一张图片
        slide.shapes.add_picture(img_path,
                                 left_pd, top_pd,
                                 width=img_weight,
                                 height=img_height)
    # 保存ppt
    prs.save(f'{os.path.join(directory, "..")}/{os.path.basename(directory)}.pptx')
    print('生成成功,文件已生成至图片的上一级目录')
